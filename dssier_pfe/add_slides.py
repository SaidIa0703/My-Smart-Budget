"""Add slides 17-50 to soutenance_CDA_MySmartBudget.pptx"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = "/Users/ahbeuhdellsaidi/Downloads/cda-folder-main/dssier_pfe/soutenance_CDA_MySmartBudget.pptx"

# Colors
C_BG       = RGBColor(0x0F, 0x17, 0x2A)
C_INDIGO   = RGBColor(0x63, 0x66, 0xF1)
C_PURPLE   = RGBColor(0xA8, 0x55, 0xF7)
C_GREEN    = RGBColor(0x34, 0xD3, 0x99)
C_RED      = RGBColor(0xF8, 0x71, 0x71)
C_YELLOW   = RGBColor(0xFF, 0xD7, 0x00)
C_DARK_BOX = RGBColor(0x1A, 0x24, 0x3E)
C_MID      = RGBColor(0x31, 0x3B, 0x6E)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0xC0, 0xC8, 0xE0)

# Font sizes
FS_LABEL   = Pt(11)
FS_TITLE   = Pt(36)
FS_SUB     = Pt(16)
FS_BOXTIT  = Pt(15)
FS_BODY    = Pt(12)
FS_SMALL   = Pt(10)
FS_TAG     = Pt(11)

# Slide dimensions (EMU)
SW = Emu(12188952)
SH = Emu(6858000)

prs = Presentation(PPTX_PATH)
blank_layout = prs.slide_layouts[6]  # Blank


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, l, t, w, h, text, bold=False, size=FS_BODY, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = size
    run.font.color.rgb = color
    return txBox


def base_slide(slide, section_label, title, subtitle, accent_color=None):
    """Draw the common frame: bg, left bar, section badge, title, divider, subtitle."""
    if accent_color is None:
        accent_color = C_INDIGO
    # Background
    add_rect(slide, 0, 0, SW, SH, C_BG)
    # Left accent bar
    add_rect(slide, 0, 0, Emu(228600), SH, accent_color)
    # Section badge
    add_rect(slide, Emu(411480), Emu(228600), Emu(2926080), Emu(384048), C_MID)
    add_text(slide, Emu(487680), Emu(254000), Emu(2834640), Emu(365760),
             section_label, bold=True, size=FS_LABEL, color=accent_color)
    # Main title
    add_text(slide, Emu(411480), Emu(731520), Emu(11338560), Emu(1005840),
             title, bold=True, size=FS_TITLE, color=C_WHITE)
    # Divider
    add_rect(slide, Emu(411480), Emu(1600200), Emu(11430000), Emu(38100), accent_color)
    # Subtitle
    add_text(slide, Emu(411480), Emu(1691640), Emu(11338560), Emu(640080),
             subtitle, bold=False, size=FS_SUB, color=C_GRAY)


def content_area_top():
    return Emu(1874519)


def add_box(slide, l, t, w, h, title, title_color, items, bullet_color=None):
    """Dark box with title and bullet items."""
    add_rect(slide, l, t, w, h, C_DARK_BOX)
    add_text(slide, Emu(l + 182880), Emu(t + 91440), Emu(w - 365760), Emu(411480),
             title, bold=True, size=FS_BOXTIT, color=title_color)
    if bullet_color is None:
        bullet_color = title_color
    item_top = t + Emu(521208)
    item_h = Emu(350520)
    for item in items:
        # bullet square
        add_rect(slide, Emu(l + 182880), Emu(item_top + 60960), Emu(182880), Emu(182880), bullet_color)
        add_text(slide, Emu(l + 457200), Emu(item_top), Emu(w - 640080), item_h,
                 item, bold=False, size=FS_BODY, color=C_GRAY)
        item_top = Emu(item_top + item_h)


def add_two_col_boxes(slide, left_title, left_color, left_items,
                      right_title, right_color, right_items):
    ct = content_area_top()
    ch = Emu(SH - ct - Emu(182880))
    col_w = Emu(5486400)
    gap = Emu(274320)
    l1 = Emu(502920)
    l2 = Emu(l1 + col_w + gap)
    add_box(slide, l1, ct, col_w, ch, left_title, left_color, left_items)
    add_box(slide, l2, ct, col_w, ch, right_title, right_color, right_items)


def add_full_box(slide, title, title_color, items, bullet_color=None):
    ct = content_area_top()
    ch = Emu(SH - ct - Emu(182880))
    add_box(slide, Emu(502920), ct, Emu(11247120), ch, title, title_color, items, bullet_color)


def add_code_slide(slide, title, title_color, code_lines, extra_items=None):
    """Slide with a code block on the left and optional bullets on the right."""
    ct = content_area_top()
    ch = Emu(SH - ct - Emu(182880))
    cw = Emu(11247120)
    add_rect(slide, Emu(502920), ct, cw, ch, C_DARK_BOX)
    add_text(slide, Emu(685800), Emu(ct + 91440), Emu(cw - 365760), Emu(411480),
             title, bold=True, size=FS_BOXTIT, color=title_color)
    code_text = "\n".join(code_lines)
    code_box = slide.shapes.add_textbox(Emu(685800), Emu(ct + 548640),
                                        Emu(cw - 365760), Emu(ch - 640080))
    code_box.text_frame.word_wrap = True
    p = code_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.bold = False
    run.font.size = Pt(10)
    run.font.color.rgb = C_GRAY


# ─── SLIDES DATA ────────────────────────────────────────────────────────────

slides_data = [

    # 17 — Modèle de données PostgreSQL
    ("s17", "04 — CONCEPTION", "Modèle de données : PostgreSQL",
     "4 tables relationnelles — clés étrangères sur user_id",
     C_INDIGO, "full",
     "Structure des tables", C_INDIGO, [
         "users — id, email, password_hash, profile_type, created_at",
         "transactions — id, user_id (FK), name, amount, category, date, is_recurring",
         "budgets — id, user_id (FK), category, limit_amount, month",
         "savings_goals — id, user_id (FK), name, target_amount, current_amount, deadline",
         "Contraintes FK : suppression cascade sur user_id",
         "Index sur user_id dans chaque table (performances requêtes)",
         "Montants négatifs = dépenses, positifs = revenus",
     ]),

    # 18 — Modèle de données MongoDB
    ("s18", "04 — CONCEPTION", "Modèle de données : MongoDB",
     "Snapshots analytiques — document auto-contenu par rapport",
     C_PURPLE, "full",
     "Document rapport (Mongoose Schema)", C_PURPLE, [
         'userId: String  — référence utilisateur PostgreSQL',
         'month: String   — ex: "2025-03"',
         'totaux: { revenus, depenses, solde }',
         'parCategorie: [{ category, total }]',
         'evolution: [{ month, revenus, depenses }]',
         'createdAt: Date — timestamp archivage',
         'Schéma flexible : champs optionnels selon version de l\'app',
     ]),

    # 19 — Pourquoi deux BDD
    ("s19", "04 — CONCEPTION", "Polyglot Persistence",
     "Chaque base fait ce qu'elle fait le mieux",
     C_GREEN, "two",
     "PostgreSQL — données transactionnelles", C_INDIGO,
     ["ACID — atomicité garantie", "Relations strictes (FK)", "Calculs agrégés SQL",
      "Données temps réel", "4 tables liées"],
     "MongoDB — données analytiques", C_PURPLE,
     ["Schéma flexible / évolutif", "Documents auto-contenus",
      "Archivage snapshots historiques", "Lecture rapide (pas de JOIN)",
      "API Mongoose simple"]),

    # 20 — Architecture API REST
    ("s20", "04 — CONCEPTION", "API REST — Organisation par domaine",
     "Chaque domaine métier isolé dans son propre fichier de routes",
     C_INDIGO, "full",
     "Endpoints REST", C_INDIGO, [
         "POST /api/auth/register  — Inscription (public)",
         "POST /api/auth/login     — Connexion, retourne JWT (public)",
         "GET/POST/PUT/DELETE /api/transactions  — JWT requis",
         "GET/POST/PUT/DELETE /api/budgets       — JWT requis",
         "GET/POST/PUT/DELETE /api/goals         — JWT requis",
         "GET /api/reports/generate — Calcul PostgreSQL, JWT requis",
         "POST/DELETE /api/reports  — Archivage MongoDB, JWT requis",
     ]),

    # 21 — Middleware JWT
    ("s21", "04 — CONCEPTION", "Authentification — Middleware JWT",
     "Vérification stateless à chaque requête protégée",
     C_INDIGO, "full",
     "Flux d'authentification", C_INDIGO, [
         "1. Client envoie: Authorization: Bearer <token>",
         "2. Middleware extrait et vérifie la signature JWT",
         "3. Si invalide → 401 Unauthorized immédiat",
         "4. Si valide → req.user = { id, email } injecté",
         "5. Route accède à req.user.id sans paramètre URL",
         "Stateless : aucun appel base de données pour vérifier",
         "Clé secrète stockée en variable d'environnement (JWT_SECRET)",
     ]),

    # 22 — Protection IDOR
    ("s22", "04 — CONCEPTION", "Sécurité — Protection IDOR",
     "Insecure Direct Object Reference — OWASP A01",
     C_RED, "two",
     "Sans protection IDOR", C_RED,
     ["GET /api/budgets/42 retourne budget de l'user 42",
      "N'importe quel user connecté peut accéder",
      "Suffit de changer l'ID dans l'URL",
      "Faille critique sur toutes les ressources"],
     "Avec ownerGuard", C_GREEN,
     ["On charge la ressource depuis la DB",
      "On compare resource.user_id === req.user.id",
      "Si différent → 403 Forbidden immédiat",
      "Appliqué sur TOUS les PUT et DELETE"]),

    # 23 — Calcul budget SQL
    ("s23", "04 — CONCEPTION", "Budget — Calcul temps réel (SQL)",
     "LEFT JOIN pour calculer les dépenses du mois courant en une requête",
     C_GREEN, "full",
     "Requête SQL optimisée", C_GREEN, [
         "SELECT b.*, COALESCE(SUM(CASE WHEN ...), 0) AS current_spending",
         "FROM budgets b LEFT JOIN transactions t",
         "  ON t.category = b.category AND t.user_id = b.user_id",
         "WHERE b.user_id = $1  GROUP BY b.id",
         "── amount < 0 filtre uniquement les dépenses",
         "── date_trunc('month', ...) limite au mois courant",
         "── COALESCE gère le cas zéro transaction (NULL → 0)",
     ]),

    # 24 — Bilan conception
    ("s24", "04 — CONCEPTION", "Bilan — Ce que la conception a changé",
     "Des décisions documentées, pas des choix au hasard",
     C_INDIGO, "two",
     "Avant conception", C_RED,
     ['"Une base de données"', "Routes en vrac dans index.js",
      "Auth vérifiée dans chaque route manuellement",
      "Pas de protection accès inter-utilisateurs",
      "Calculs faits en JavaScript après 2 requêtes"],
     "Après conception", C_GREEN,
     ["PostgreSQL + MongoDB (polyglot)", "Domaines métier séparés (5 fichiers routes)",
      "Middleware JWT centralisé réutilisable",
      "ownerGuard sur tous les PUT/DELETE",
      "JOIN SQL : 1 requête, 0 traitement JS"]),

    # 25 — Plan démonstration
    ("s25", "05 — DÉMO", "Démonstration — Parcours Utilisateur",
     "Application en conditions réelles — Docker local",
     C_GREEN, "full",
     "Étapes de la démo", C_GREEN, [
         "1. Inscription + Questionnaire profil (Salarié)",
         "2. Dashboard personnalisé — KPIs temps réel",
         "3. Ajout de transactions (revenus + dépenses)",
         "4. Gestion des budgets + alertes dépassement",
         "5. Objectifs d'épargne — barre de progression",
         "6. Rapports analytiques + archivage MongoDB",
         "7. Export CSV côté client",
     ]),

    # 26 — Demo Inscription Questionnaire
    ("s26", "05 — DÉMO", "Onboarding — Questionnaire Adaptatif",
     "3 profils — questions différentes selon le choix",
     C_GREEN, "two",
     "Profils disponibles", C_GREEN,
     ["🎓 Étudiant — bourses, aide familiale, jobs étudiants",
      "💼 Salarié — salaire net, primes, charges",
      "🏢 Entrepreneur — CA, charges professionnelles"],
     "Ce que le questionnaire configure", C_INDIGO,
     ["Catégories par défaut dans le dashboard",
      "KPIs adaptés au profil sélectionné",
      "Suggestions de budgets selon les revenus",
      "Labels personnalisés (ex: 'Bourse' vs 'Salaire')"]),

    # 27 — Demo Dashboard
    ("s27", "05 — DÉMO", "Dashboard — Vue Temps Réel",
     "KPIs calculés depuis les transactions du mois courant",
     C_GREEN, "full",
     "Composants du dashboard", C_GREEN, [
         "Solde = Revenus − Dépenses (calcul côté client)",
         "Liste des dernières transactions (5 plus récentes)",
         "Graphique répartition dépenses par catégorie",
         "Badge rouge dans navbar si budget dépassé",
         "Bouton d'accès rapide à chaque section",
         "Nom de l'utilisateur affiché (depuis le token JWT)",
         "3 variantes : Étudiant / Salarié / Entrepreneur",
     ]),

    # 28 — Demo Budgets Alertes
    ("s28", "05 — DÉMO", "Budgets — Alertes Dynamiques",
     "Calcul temps réel : LEFT JOIN PostgreSQL",
     C_YELLOW, "two",
     "Indicateurs visuels", C_YELLOW,
     ["Barre de progression par catégorie",
      "Vert : dépenses < 80% du budget",
      "Orange : dépenses entre 80% et 100%",
      "Rouge : budget dépassé (> 100%)",
      "Bannière d'alerte en haut de page"],
     "CRUD Budgets", C_INDIGO,
     ["Créer un budget par catégorie + montant limite",
      "Modifier le montant limite",
      "Supprimer un budget (confirmation)",
      "409 si catégorie déjà utilisée",
      "Badge navbar mis à jour automatiquement"]),

    # 29 — Demo Rapports Export
    ("s29", "05 — DÉMO", "Rapports — Analytics + Archivage MongoDB",
     "Calcul PostgreSQL → affichage → archivage optionnel MongoDB",
     C_PURPLE, "full",
     "Fonctionnalités de la page Rapports", C_PURPLE, [
         "Sélecteur de mois — données filtrées dynamiquement",
         "4 KPIs : Revenus totaux / Dépenses / Solde / Économies",
         "Graphique barres : dépenses par catégorie",
         "Courbe évolution : 6 mois glissants",
         'Bouton "Archiver dans MongoDB" — snapshot permanent',
         "Liste des rapports archivés avec date",
         "Export CSV : fichier généré côté client (0 serveur)",
     ]),

    # 30 — Frontend Stack
    ("s30", "06 — FRONTEND", "Frontend — Next.js 15 + React 19",
     "App Router, TypeScript strict, Tailwind CSS",
     C_INDIGO, "two",
     "Arborescence du projet", C_INDIGO,
     ["src/app/         — Pages (App Router)",
      "components/      — Composants réutilisables",
      "src/types/       — Interfaces TypeScript",
      "src/lib/         — storage.ts, hooks custom",
      "public/          — Assets statiques"],
     "Pages de l'application", C_PURPLE,
     ["/login  /register  /questionnaire",
      "/dashboard  (3 variantes profil)",
      "/transaction  — CRUD transactions",
      "/budgets  — CRUD budgets + alertes",
      "/goals  — CRUD objectifs épargne",
      "/reports  — Analytics + archivage"]),

    # 31 — Frontend Routing
    ("s31", "06 — FRONTEND", "App Router — Navigation & Guards",
     "Chaque page vérifie le JWT au montage",
     C_INDIGO, "full",
     "Protection des routes", C_INDIGO, [
         "useEffect au montage : lecture sessionStorage || localStorage",
         "Si token absent → router.push('/login') immédiat",
         "Si token présent → fetch des données avec Authorization header",
         "Navbar affichée uniquement pour les pages protégées",
         "Redirection post-login vers la page demandée",
         "Logout : clear sessionStorage ET localStorage + redirect /login",
         "Badge budgets : recalculé à chaque changement de pathname",
     ]),

    # 32 — Hook partagé
    ("s32", "06 — FRONTEND", "useProfileDashboard — Hook Partagé",
     "1 hook custom → 3 dashboards — zéro duplication",
     C_PURPLE, "two",
     "Problème sans hook", C_RED,
     ["3 dashboards, même logique de fetch",
      "150 lignes dupliquées × 3 = 450 lignes",
      "Un bug corrigé dans 1 = à corriger dans 3",
      "Incohérences entre profils inévitables"],
     "Solution : hook custom", C_GREEN,
     ["useProfileDashboard() → { transactions, loading, kpis }",
      "Fetch /api/transactions avec JWT",
      "Calcul KPIs (revenus, dépenses, solde)",
      "Chaque dashboard = rendu uniquement",
      "Bug corrigé en 1 seul endroit"]),

    # 33 — Sessions Storage
    ("s33", "06 — FRONTEND", "Sécurité Frontend — Gestion des Sessions",
     "sessionStorage par défaut — localStorage en fallback",
     C_RED, "two",
     "sessionStorage (défaut)", C_GREEN,
     ["Effacé à la fermeture du navigateur",
      "Accessible uniquement onglet courant",
      "JWT stocké ici par saveSession()",
      "Protège contre vol de session cross-tab",
      "XSS limité : scripts injectés lisent localStorage"],
     "Règles appliquées", C_INDIGO,
     ["Lecture : sessionStorage || localStorage (fallback)",
      "Écriture (login) : sessionStorage uniquement",
      "Logout : clear() sur les DEUX storages",
      "Token envoyé en header Authorization: Bearer",
      "Jamais de token dans les paramètres URL"]),

    # 34 — TypeScript
    ("s34", "06 — FRONTEND", "TypeScript — Contrats entre Couches",
     "Les interfaces documentent et protègent les échanges frontend/backend",
     C_PURPLE, "full",
     "Interfaces métier principales", C_PURPLE, [
         "Transaction : { id, name, amount, category, date, is_recurring? }",
         "Budget : { id, category, limit_amount, current_spending }",
         "SavingsGoal : { id, name, target_amount, current_amount, deadline? }",
         "Report : { userId, month, totaux, parCategorie, evolution }",
         "Props typées sur chaque composant React",
         "TypeScript strict : aucun 'any' implicite toléré",
         "Erreur de compilation si backend change un nom de champ",
     ]),

    # 35 — Export CSV
    ("s35", "06 — FRONTEND", "Export CSV — Zéro Dépendance Serveur",
     "Blob + URL.createObjectURL — génération entièrement côté client",
     C_GREEN, "full",
     "Implémentation export CSV", C_GREEN, [
         "Header CSV : 'Date,Description,Catégorie,Montant'",
         "Chaque transaction → ligne CSV formatée",
         "new Blob([csv], { type: 'text/csv' })",
         "URL.createObjectURL(blob) → URL temporaire",
         "Clic programmatique sur <a download='transactions.csv'>",
         "URL.revokeObjectURL() → nettoyage mémoire",
         "Aucun round-trip serveur, aucune dépendance externe",
     ]),

    # 36 — Backend Structure
    ("s36", "07 — BACKEND", "Backend — Architecture Express Modulaire",
     "Separation of Concerns — 1 fichier par domaine métier",
     C_INDIGO, "two",
     "Arborescence backend", C_INDIGO,
     ["index.js         — Entry point + middlewares",
      "routes/auth.js   — /api/auth/*",
      "routes/transactions.js",
      "routes/budgets.js",
      "routes/goals.js",
      "routes/reports.js",
      "middleware/auth.js — JWT verify",
      "models/report.js  — Mongoose schema",
      "db.js            — pg-promise pool"],
     "Principe d'extension", C_GREEN,
     ["Ajouter un domaine = 1 fichier routes/X.js",
      "Importer en 1 ligne dans index.js",
      "Middleware JWT déjà appliqué automatiquement",
      "Tests isolés par fichier de routes",
      "Aucun couplage entre domaines"]),

    # 37 — Backend DB connections
    ("s37", "07 — BACKEND", "Connexions Base de Données",
     "pg-promise (PostgreSQL) + Mongoose (MongoDB) — via env vars",
     C_INDIGO, "full",
     "Configuration des connexions", C_INDIGO, [
         "DATABASE_URL → pg-promise pool (connexions réutilisées)",
         "MONGODB_URI  → mongoose.connect() avec reconnexion auto",
         "Aucun credential en dur dans le code source",
         ".env dans .gitignore — jamais committé",
         "docker-compose.yml injecte les variables au démarrage",
         "Pool PostgreSQL : max 10 connexions simultanées",
         "Mongoose : retry automatique si MongoDB redémarre",
     ]),

    # 38 — Backend route budgets
    ("s38", "07 — BACKEND", "Route Complexe — GET /api/budgets",
     "Un seul SQL avec LEFT JOIN remplace 2 requêtes + fusion JS",
     C_GREEN, "full",
     "Approche et bénéfices", C_GREEN, [
         "SELECT b.*, COALESCE(SUM(CASE WHEN ...), 0) AS current_spending",
         "LEFT JOIN transactions ON category + user_id",
         "CASE WHEN amount < 0 AND date_trunc('month') = NOW()",
         "GROUP BY b.id — 1 ligne par budget avec total dépenses",
         "Avantage : 1 round-trip DB au lieu de 1+N",
         "COALESCE : budget sans transaction → 0 (pas NULL)",
         "Résultat prêt à l'emploi pour le frontend",
     ]),

    # 39 — Backend reports generate
    ("s39", "07 — BACKEND", "Génération de Rapports — 3 Requêtes SQL",
     "Promise.all — exécution parallèle des 3 agrégations",
     C_PURPLE, "two",
     "3 dimensions analytiques", C_PURPLE,
     ["1. Totaux : SUM revenus / SUM dépenses du mois",
      "2. Par catégorie : GROUP BY category, SUM(ABS(amount))",
      "3. Évolution 6 mois : GROUP BY date_trunc('month')",
      "Promise.all([q1, q2, q3]) — parallèle"],
     "Flux complet", C_INDIGO,
     ["GET /api/reports/generate → 3 SQL en parallèle",
      "Assemblage JSON : { totaux, parCategorie, evolution }",
      "Frontend reçoit et affiche les graphiques",
      'Bouton "Archiver" → POST /api/reports (MongoDB)',
      "GET /api/reports → liste snapshots archivés"]),

    # 40 — Variables env
    ("s40", "07 — BACKEND", "Configuration — 12-Factor App",
     "La config dans l'environnement, jamais dans le code",
     C_YELLOW, "full",
     "Variables d'environnement", C_YELLOW, [
         "DATABASE_URL  — postgresql://user:pass@postgres:5432/mydb",
         "MONGODB_URI   — mongodb://mongo:27017/mybudget",
         "JWT_SECRET    — chaîne aléatoire longue (256 bits)",
         "PORT          — 5001 (backend Express)",
         "FRONTEND_URL  — http://localhost:3000 (CORS whitelist)",
         ".env.example fourni dans le repo (sans valeurs réelles)",
         "En production Docker : variables dans docker-compose.yml",
     ]),

    # 41 — Middlewares sécurité
    ("s41", "07 — BACKEND", "Sécurité Backend — Stack Middlewares",
     "4 couches de protection avant d'atteindre les routes",
     C_RED, "two",
     "Middlewares appliqués", C_RED,
     ["helmet() — 15+ headers HTTP sécurisés",
      "cors({ origin: FRONTEND_URL }) — whitelist stricte",
      "rateLimit(100 req / 15min / IP) — anti brute-force",
      "express.json() — parse body JSON uniquement"],
     "Ce que chaque middleware bloque", C_INDIGO,
     ["Helmet : clickjacking, MIME sniffing, XSS basique",
      "CORS : requêtes cross-origin non autorisées",
      "Rate limit : attaques par force brute sur /auth/login",
      "JSON parse : payload non-JSON rejeté (400)",
      "JWT middleware : toute requête sans token valide"]),

    # 42 — OWASP
    ("s42", "08 — SÉCURITÉ", "Sécurité — Couverture OWASP Top 10",
     "6 risques critiques couverts dans ce MVP",
     C_RED, "two",
     "Risques couverts ✅", C_GREEN,
     ["A01 Broken Access Control → ownerGuard + JWT",
      "A02 Cryptographic Failures → bcrypt salt 12 + HTTPS",
      "A03 Injection → requêtes paramétrées ($1, $2)",
      "A05 Security Misconfiguration → Helmet + CORS",
      "A07 Auth Failures → JWT stateless + rate limiting",
      "A09 Logging → Morgan + console.error"],
     "Roadmap sécurité 🔜", C_YELLOW,
     ["A04 XSS côté frontend → Content Security Policy",
      "A06 Outdated Components → Dependabot alerts",
      "A08 SSRF → validation des URLs utilisateur",
      "A10 Supply Chain → npm audit en CI",
      "Pentest → OWASP ZAP automatisé"]),

    # 43 — Bcrypt
    ("s43", "08 — SÉCURITÉ", "Authentification — bcrypt",
     "Mots de passe hachés, jamais stockés en clair",
     C_RED, "full",
     "Fonctionnement de bcrypt", C_RED, [
         "Inscription : bcrypt.hash(password, 12) → hash unique",
         "Salt factor 12 : ~250ms par hash (délibérément lent)",
         "Connexion : bcrypt.compare(password, stored_hash)",
         "Compare est résistant aux timing attacks",
         "Chaque hash est unique même pour le même mot de passe",
         "Rainbow tables inefficaces (salt aléatoire intégré)",
         "Le hash ne permet JAMAIS de retrouver le mot de passe",
     ]),

    # 44 — SQL injection
    ("s44", "08 — SÉCURITÉ", "Injection SQL — Prévention par Conception",
     "Paramètres positionnels — jamais de concaténation de chaînes",
     C_RED, "two",
     "Approche dangereuse ❌", C_RED,
     ['db.query(`SELECT * FROM users', '  WHERE email = \'${email}\'`)',
      "Si email = \"' OR '1'='1\"",
      "→ retourne TOUS les utilisateurs",
      "Si email = \"; DROP TABLE users;--\"",
      "→ supprime toute la table"],
     "Approche sécurisée ✅", C_GREEN,
     ["db.query('SELECT * FROM users WHERE email = $1', [email])",
      "La valeur est envoyée séparément du SQL",
      "pg-promise l'échappe automatiquement",
      "Aucune interprétation SQL de la valeur",
      "Appliqué sur 100% des requêtes du projet"]),

    # 45 — SonarCloud
    ("s45", "08 — SÉCURITÉ", "Qualité & Sécurité — SonarCloud",
     "Analyse statique à chaque push — Quality Gate bloque les PRs",
     C_YELLOW, "full",
     "Métriques SonarCloud", C_YELLOW, [
         "0 Security Hotspot non résolu",
         "0 Bug critique ou bloquant",
         "Coverage : > 80% sur les routes backend",
         "Code Smells : tous résolus avant merge",
         "Quality Gate : statut affiché dans le README",
         "Intégration GitHub Actions : scan automatique sur push",
         "Rapport de couverture lcov généré par Jest",
     ]),

    # 46 — Tests stratégie
    ("s46", "09 — TESTS", "Tests — Stratégie et Couverture",
     "43 tests, 5 suites — Jest + Supertest",
     C_YELLOW, "two",
     "Architecture des tests", C_YELLOW,
     ["auth.routes.test.js     — 8 tests",
      "transactions.routes.test.js — 7 tests",
      "budgets.routes.test.js  — 10 tests",
      "goals.routes.test.js    — 10 tests",
      "reports.routes.test.js  — 8 tests",
      "Total : 43 tests, 5 suites"],
     "Mocks utilisés", C_INDIGO,
     ["jest.mock('../db') → faux pg-promise",
      "jest.mock('../models/report') → faux Mongoose",
      "mockDb.query.mockResolvedValueOnce([...data...])",
      "mockReport.find → { sort: jest.fn() } (chaîne fluente)",
      "Token JWT généré en setup pour les routes protégées",
      "Tests isolés : aucun état partagé entre suites"]),

    # 47 — Test exemple
    ("s47", "09 — TESTS", "Tests — Exemple Concret (budgets)",
     "Cas de test : 409 si catégorie déjà existante",
     C_YELLOW, "full",
     "Ce que ce test vérifie", C_YELLOW, [
         "mockDb.oneOrNone retourne un budget existant (mock)",
         "POST /api/budgets avec category: 'Alimentation'",
         "Authorization: Bearer <token valide> dans le header",
         "Résultat attendu : status 409",
         "Body attendu : { error: '...existe déjà...' }",
         "Sans mock : dépendrait d'une DB en état connu → fragile",
         "Avec mock : déterministe, rapide, reproductible",
     ]),

    # 48 — GitHub Actions
    ("s48", "09 — TESTS", "CI/CD — GitHub Actions Pipeline",
     "Déclenchement automatique sur push et pull_request",
     C_GREEN, "two",
     "Étapes du pipeline", C_GREEN,
     ["1. Checkout du code (actions/checkout)",
      "2. Setup Node.js 20 (actions/setup-node)",
      "3. npm install (installation dépendances)",
      "4. npm test (43 tests Jest)",
      "5. Génération rapport coverage lcov",
      "6. SonarCloud Scan (quality gate)"],
     "Effets sur le workflow", C_INDIGO,
     ["PR bloquée si un test échoue",
      "PR bloquée si Quality Gate SonarCloud KO",
      "Badge vert/rouge dans le README du repo",
      "Historique des runs visible sur GitHub",
      "Notification automatique si build casse"]),

    # 49 — Docker Compose
    ("s49", "10 — DEVOPS", "Déploiement — Docker Compose 4 Services",
     "'docker compose up' démarre toute l'infrastructure",
     C_GREEN, "two",
     "Les 4 services", C_GREEN,
     ["nginx  — Reverse proxy :80 → frontend/backend",
      "frontend — Next.js :3000",
      "backend  — Express :5001",
      "postgres — PostgreSQL :5432 (volume persistant)",
      "mongo    — MongoDB :27017 (volume persistant)"],
     "Avantages Docker", C_INDIGO,
     ["Environnement identique dev / prod / jury",
      "Un seul fichier docker-compose.yml à configurer",
      "Images versionnées sur Docker Hub",
      "init.sql exécuté au premier démarrage PostgreSQL",
      "Healthchecks : backend attend que la DB soit prête"]),

    # 50 — Conclusion
    ("s50", "CONCLUSION", "Bilan & Perspectives",
     "My Smart Budget — du concept à l'application déployable",
     C_INDIGO, "two",
     "Réalisé en 8 semaines ✅", C_GREEN,
     ["Application full-stack fonctionnelle et déployable",
      "JWT + bcrypt + OWASP — sécurité production",
      "4 domaines métier (transactions, budgets, goals, reports)",
      "Polyglot persistence PostgreSQL + MongoDB",
      "43 tests automatisés + SonarCloud + CI/CD",
      "Docker Compose — 1 commande = tout démarre"],
     "Roadmap 🚀", C_PURPLE,
     ["Application mobile React Native (iOS/Android)",
      "IA : catégorisation automatique des transactions",
      "Open Banking PSD2 — import relevés bancaires",
      "Notifications push : budget dépassé en temps réel",
      "Budgets partagés (famille / colocation)",
      "Mode multi-devise avec taux de change live"]),
]


def build_slide(prs, data):
    sid, section, title, subtitle, accent, layout_type, *rest = data
    slide = prs.slides.add_slide(blank_layout)
    base_slide(slide, section, title, subtitle, accent)

    if layout_type == "full":
        box_title, box_color = rest[0], rest[1]
        items = rest[2]
        add_full_box(slide, box_title, box_color, items, box_color)

    elif layout_type == "two":
        lt, lc, li, rt, rc, ri = rest[0], rest[1], rest[2], rest[3], rest[4], rest[5]
        add_two_col_boxes(slide, lt, lc, li, rt, rc, ri)


for data in slides_data:
    build_slide(prs, data)

prs.save(PPTX_PATH)
print(f"Saved. Total slides: {len(prs.slides)}")
