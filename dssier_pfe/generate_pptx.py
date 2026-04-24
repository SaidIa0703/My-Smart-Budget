"""
Génère le PowerPoint de soutenance CDA — My Smart Budget
Style : dark mode indigo/violet (inspiré Canva)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0F, 0x17, 0x2A)   # Fond principal très sombre
INDIGO      = RGBColor(0x63, 0x66, 0xF1)   # Indigo vif
PURPLE      = RGBColor(0xA8, 0x55, 0xF7)   # Violet
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)   # Jaune Colint
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xC0, 0xC8, 0xE0)
DARK_CARD   = RGBColor(0x1A, 0x24, 0x3E)
GREEN       = RGBColor(0x34, 0xD3, 0x99)
RED_SOFT    = RGBColor(0xF8, 0x71, 0x71)
INDIGO_SOFT = RGBColor(0x31, 0x3B, 0x6E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # Totalement vide

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def bg(slide):
    """Fond sombre + bande déco gauche indigo"""
    add_rect(slide, 0, 0, W, H, BG_DARK)
    add_rect(slide, 0, 0, Inches(0.25), H, INDIGO)

def accent_bar(slide, y=Inches(1.8), color=INDIGO):
    add_rect(slide, Inches(0.45), y, Inches(12.5), Pt(3), color)

def chip(slide, label, x, y, color=INDIGO):
    w = Inches(1.6)
    h = Inches(0.38)
    add_rect(slide, x, y, w, h, color)
    add_text(slide, label, x, y + Pt(2), w, h,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def section_badge(slide, text, x=Inches(0.45), y=Inches(0.25)):
    add_rect(slide, x, y, Inches(3.2), Inches(0.42), INDIGO_SOFT)
    add_text(slide, text, x + Pt(6), y + Pt(2), Inches(3.1), Inches(0.4),
             font_size=11, bold=True, color=INDIGO, align=PP_ALIGN.LEFT)

def slide_title(slide, title, subtitle=None, badge=None):
    if badge:
        section_badge(slide, badge)
    add_text(slide, title,
             Inches(0.45), Inches(0.8), Inches(12.4), Inches(1.1),
             font_size=36, bold=True, color=WHITE)
    accent_bar(slide, Inches(1.75))
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(1.85), Inches(12.4), Inches(0.7),
                 font_size=16, color=LIGHT_GRAY)

def card(slide, x, y, w, h, title=None, body_lines=None,
         title_color=INDIGO, body_color=LIGHT_GRAY, bg_color=DARK_CARD):
    add_rect(slide, x, y, w, h, bg_color)
    cy = y + Inches(0.15)
    if title:
        add_text(slide, title, x + Inches(0.18), cy,
                 w - Inches(0.3), Inches(0.45),
                 font_size=13, bold=True, color=title_color)
        cy += Inches(0.42)
    if body_lines:
        for line in body_lines:
            add_text(slide, line, x + Inches(0.18), cy,
                     w - Inches(0.3), Inches(0.38),
                     font_size=12, color=body_color)
            cy += Inches(0.35)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
# Grande bande colorée
add_rect(s, 0, Inches(2.5), W, Inches(3.2), INDIGO_SOFT)
add_rect(s, 0, Inches(2.5), Inches(0.55), Inches(3.2), INDIGO)

# MSB logo
add_rect(s, Inches(0.55), Inches(0.3), Inches(1.05), Inches(1.05), INDIGO)
add_text(s, "MSB", Inches(0.55), Inches(0.38), Inches(1.05), Inches(0.9),
         font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Titre principal
add_text(s, "My Smart Budget",
         Inches(0.75), Inches(2.65), Inches(12), Inches(1.5),
         font_size=52, bold=True, color=WHITE)

add_text(s, "Application web de gestion budgétaire personnalisée",
         Inches(0.75), Inches(4.0), Inches(12), Inches(0.7),
         font_size=20, color=LIGHT_GRAY)

# Chips technos
chip_data = [
    ("Next.js 15", Inches(0.75)),
    ("Node.js + Express", Inches(2.45)),
    ("PostgreSQL", Inches(4.15)),
    ("MongoDB", Inches(5.85)),
    ("Docker", Inches(7.55)),
]
for label, cx in chip_data:
    chip(s, label, cx, Inches(4.85), INDIGO)

# Infos candidat
add_rect(s, 0, Inches(6.3), W, Inches(1.2), RGBColor(0x0B, 0x10, 0x1F))
add_text(s, "Titre CDA — Concepteur Développeur d'Applications  |  Colint School  |  2025",
         Inches(0.75), Inches(6.5), Inches(12), Inches(0.6),
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Plan
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Plan de présentation", "~60 minutes", badge="SOMMAIRE")

plan = [
    ("01", "Introduction en anglais (QQOQCP)",    "5 min",  INDIGO),
    ("02", "Contexte & problématique",             "5 min",  INDIGO),
    ("03", "Méthodologie & organisation",          "5 min",  PURPLE),
    ("04", "Conception : architecture, BDD, UML", "10 min", PURPLE),
    ("05", "Démonstration live",                  "8 min",  GREEN),
    ("06", "Code Frontend — points clés",          "7 min",  INDIGO),
    ("07", "Code Backend & API REST",              "7 min",  INDIGO),
    ("08", "Sécurité applicative",                "5 min",  RED_SOFT),
    ("09", "Tests & qualité SonarCloud",           "4 min",  YELLOW),
    ("10", "Déploiement Docker & DevOps",          "4 min",  GREEN),
]

cols = [plan[:5], plan[5:]]
for col_idx, col in enumerate(cols):
    x_base = Inches(0.55) + col_idx * Inches(6.5)
    for i, (num, label, dur, col_color) in enumerate(col):
        y = Inches(2.1) + i * Inches(0.95)
        add_rect(s, x_base, y, Inches(6.2), Inches(0.82), DARK_CARD)
        add_rect(s, x_base, y, Inches(0.55), Inches(0.82), col_color)
        add_text(s, num, x_base, y + Pt(8), Inches(0.55), Inches(0.5),
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, label,
                 x_base + Inches(0.65), y + Pt(6), Inches(4.5), Inches(0.5),
                 font_size=13, bold=True, color=WHITE)
        add_text(s, dur,
                 x_base + Inches(5.15), y + Pt(6), Inches(0.9), Inches(0.5),
                 font_size=12, color=col_color, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction anglais
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Introduction — QQOQCP",
            "Présentation en anglais • 5 minutes", badge="01 — ENGLISH")

qqoqcp = [
    ("WHO?",   "Students, employees and families who want to manage their personal finances"),
    ("WHAT?",  "A full-stack web app to track income, expenses, budgets, goals & reports"),
    ("WHERE?", "Accessible via web browser — deployed online with Docker containers"),
    ("WHEN?",  "Developed over ~6 months during my work-study program (2024–2025)"),
    ("HOW?",   "Next.js + Node.js/Express + PostgreSQL + MongoDB + Docker"),
    ("WHY?",   "People struggle without the right tools — MSB provides a smart, secure, personalized experience"),
]

cols2 = [qqoqcp[:3], qqoqcp[3:]]
for ci, col in enumerate(cols2):
    for ri, (q, a) in enumerate(col):
        x = Inches(0.55) + ci * Inches(6.5)
        y = Inches(2.1) + ri * Inches(1.55)
        add_rect(s, x, y, Inches(6.2), Inches(1.4), DARK_CARD)
        add_rect(s, x, y, Inches(1.3), Inches(1.4), INDIGO)
        add_text(s, q, x, y + Inches(0.4), Inches(1.3), Inches(0.7),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, a, x + Inches(1.45), y + Inches(0.2),
                 Inches(4.55), Inches(1.0),
                 font_size=12, color=LIGHT_GRAY, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Contexte & Fonctionnalités
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Contexte & Fonctionnalités",
            "Une application adaptée à chaque profil utilisateur", badge="02 — PROJET")

# Problématique
add_rect(s, Inches(0.55), Inches(2.05), Inches(5.9), Inches(1.05), INDIGO_SOFT)
add_text(s, "Problématique",
         Inches(0.75), Inches(2.1), Inches(5.7), Inches(0.4),
         font_size=12, bold=True, color=INDIGO)
add_text(s, "Comment aider les particuliers à mieux gérer leur budget "
            "avec une expérience personnalisée selon leur profil ?",
         Inches(0.75), Inches(2.45), Inches(5.7), Inches(0.6),
         font_size=13, color=WHITE, wrap=True)

# 3 profils
for i, (icon, label, desc) in enumerate([
    ("🎓", "Étudiant",  "Budget serré, revenus variables"),
    ("💼", "Salarié",   "Revenus fixes, projets long terme"),
    ("👨‍👩‍👧", "Famille", "Multi-personnes, dépenses partagées"),
]):
    x = Inches(0.55) + i * Inches(2.0)
    add_rect(s, x, Inches(3.3), Inches(1.85), Inches(1.1), DARK_CARD)
    add_text(s, icon,   x, Inches(3.33), Inches(1.85), Inches(0.5),
             font_size=18, align=PP_ALIGN.CENTER)
    add_text(s, label,  x, Inches(3.7),  Inches(1.85), Inches(0.38),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc,   x, Inches(4.0),  Inches(1.85), Inches(0.38),
             font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Features liste
features = [
    "Inscription / connexion sécurisée (JWT + bcrypt)",
    "Questionnaire d'onboarding personnalisé (7 étapes)",
    "Dashboard avec statistiques temps réel",
    "CRUD Transactions avec protection IDOR",
    "Gestion des budgets par catégorie",
    "Objectifs d'épargne (Goals)",
    "Rapports analytiques (MongoDB)",
    "Interface responsive mobile + desktop",
]
add_rect(s, Inches(6.65), Inches(2.05), Inches(6.2), Inches(5.2), DARK_CARD)
add_text(s, "Fonctionnalités",
         Inches(6.85), Inches(2.15), Inches(6.0), Inches(0.45),
         font_size=14, bold=True, color=INDIGO)
for i, f in enumerate(features):
    add_rect(s, Inches(6.85), Inches(2.6) + i * Inches(0.57), Inches(0.28), Inches(0.28), GREEN)
    add_text(s, f,
             Inches(7.25), Inches(2.58) + i * Inches(0.57), Inches(5.4), Inches(0.42),
             font_size=12, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Méthodologie
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Méthodologie & Organisation",
            "Approche Agile adaptée — outils professionnels", badge="03 — MÉTHODE")

# Agile
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.9), DARK_CARD)
add_text(s, "Méthode Agile / Scrum Solo",
         Inches(0.75), Inches(2.15), Inches(5.8), Inches(0.5),
         font_size=15, bold=True, color=INDIGO)
agile_items = [
    ("Sprints",       "1 à 2 semaines par fonctionnalité"),
    ("Backlog",       "Items priorisés par valeur métier"),
    ("Git Flow",      "Branches feat/, fix/, docs/ + main"),
    ("Commits",       "Conventionnels : feat, fix, refactor, docs"),
    ("Code review",   "SonarCloud à chaque push"),
    ("Retrospective", "Ajustement de l'architecture si nécessaire"),
]
for i, (t, d) in enumerate(agile_items):
    y = Inches(2.75) + i * Inches(0.67)
    add_rect(s, Inches(0.75), y, Inches(1.5), Inches(0.5), INDIGO_SOFT)
    add_text(s, t, Inches(0.75), y + Pt(4), Inches(1.5), Inches(0.42),
             font_size=11, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    add_text(s, d, Inches(2.4), y + Pt(4), Inches(4.0), Inches(0.42),
             font_size=12, color=LIGHT_GRAY)

# Outils
add_rect(s, Inches(6.8), Inches(2.05), Inches(6.0), Inches(4.9), DARK_CARD)
add_text(s, "Stack & Outils",
         Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5),
         font_size=15, bold=True, color=PURPLE)
tools = [
    ("Frontend",    "Next.js 15 · React 19 · TypeScript · Tailwind CSS"),
    ("Backend",     "Node.js · Express · REST API"),
    ("BDD SQL",     "PostgreSQL 16 · pg-promise · indexes"),
    ("BDD NoSQL",   "MongoDB 7 · Mongoose ODM"),
    ("Sécurité",    "JWT · bcrypt · Helmet · Rate Limit"),
    ("Tests",       "Jest · Supertest · coverage lcov"),
    ("DevOps",      "Docker · Docker Compose · Docker Hub"),
    ("Qualité",     "SonarCloud · ESLint · TypeScript strict"),
    ("IDE & VCS",   "VS Code · GitHub"),
]
for i, (t, d) in enumerate(tools):
    y = Inches(2.7) + i * Inches(0.49)
    add_text(s, f"▸  {t}:", Inches(7.0), y, Inches(1.5), Inches(0.4),
             font_size=11, bold=True, color=PURPLE)
    add_text(s, d, Inches(8.4), y, Inches(4.2), Inches(0.4),
             font_size=11, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture 3-tiers
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Architecture 3 Tiers",
            "Séparation stricte des responsabilités", badge="04 — ARCHITECTURE")

layers = [
    (INDIGO,  "COUCHE PRÉSENTATION",   "Next.js 15 + React 19 + TypeScript",
              "Tailwind CSS · Recharts · App Router · Composants réutilisables"),
    (PURPLE,  "COUCHE MÉTIER (API REST)", "Node.js + Express",
              "Middleware JWT · Helmet · Rate Limiting · Validation · ownerGuard"),
    (GREEN,   "COUCHE ACCÈS DONNÉES",  "pg-promise (PostgreSQL) · Mongoose (MongoDB)",
              "Requêtes paramétrées · Indexes · Upsert · Agrégations SQL"),
    (YELLOW,  "COUCHE PERSISTANCE",    "PostgreSQL 16 · MongoDB 7",
              "ACID · Intégrité référentielle · Schéma flexible NoSQL"),
]

y_start = Inches(2.05)
layer_h = Inches(1.13)
arrow_h = Inches(0.2)

for i, (col, title, tech, detail) in enumerate(layers):
    y = y_start + i * (layer_h + arrow_h)
    add_rect(s, Inches(0.55), y, Inches(12.3), layer_h, DARK_CARD)
    add_rect(s, Inches(0.55), y, Inches(0.4), layer_h, col)
    add_text(s, title,
             Inches(1.15), y + Pt(8), Inches(3.8), Inches(0.45),
             font_size=12, bold=True, color=col)
    add_text(s, tech,
             Inches(5.0), y + Pt(8), Inches(7.6), Inches(0.45),
             font_size=13, bold=True, color=WHITE)
    add_text(s, detail,
             Inches(1.15), y + Pt(40), Inches(11.5), Inches(0.5),
             font_size=11, color=LIGHT_GRAY)
    # Arrow
    if i < 3:
        ay = y + layer_h
        add_text(s, "▼", Inches(6.4), ay, Inches(0.6), Inches(0.22),
                 font_size=12, color=INDIGO, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Conception BDD
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Conception Base de Données",
            "PostgreSQL (relationnel) + MongoDB (NoSQL)", badge="04 — BDD")

# PostgreSQL tables
tables = [
    ("users",         ["id (PK)", "name", "email (UNIQUE)", "password (bcrypt)", "created_at", "updated_at"]),
    ("transactions",  ["id (PK)", "user_id (FK→users)", "name", "category", "amount DECIMAL", "date", "is_recurring", "updated_at"]),
    ("budgets",       ["id (PK)", "user_id (FK→users)", "category", "limit DECIMAL", "created_at"]),
    ("user_profiles", ["user_id (PK/FK)", "profile_type", "revenus_mensuels", "situation_familiale", "objectifs_epargne", "categories_prioritaires"]),
]

for ti, (tname, cols) in enumerate(tables):
    x = Inches(0.55) + (ti % 2) * Inches(4.5)
    y = Inches(2.05) + (ti // 2) * Inches(2.5)
    w = Inches(4.3)
    add_rect(s, x, y, w, Inches(0.42), INDIGO)
    add_text(s, tname, x, y + Pt(3), w, Inches(0.36),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ci, col in enumerate(cols):
        row_y = y + Inches(0.42) + ci * Inches(0.32)
        row_bg = DARK_CARD if ci % 2 == 0 else RGBColor(0x14, 0x1E, 0x35)
        add_rect(s, x, row_y, w, Inches(0.32), row_bg)
        add_text(s, col, x + Pt(8), row_y + Pt(3), w - Pt(12), Inches(0.3),
                 font_size=11, color=LIGHT_GRAY)

# MongoDB
add_rect(s, Inches(9.35), Inches(2.05), Inches(3.5), Inches(5.1), DARK_CARD)
add_rect(s, Inches(9.35), Inches(2.05), Inches(3.5), Inches(0.42), GREEN)
add_text(s, "Report (MongoDB)", Inches(9.35), Inches(2.08), Inches(3.5), Inches(0.36),
         font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
mongo_fields = ["_id: ObjectId", "userId: String", "title: String",
                "description: String", "type: monthly|yearly|custom",
                'data: Mixed  { "montant": ..., }', "createdAt: ISODate", "updatedAt: ISODate"]
for i, f in enumerate(mongo_fields):
    yy = Inches(2.55) + i * Inches(0.38)
    add_text(s, f, Inches(9.55), yy, Inches(3.1), Inches(0.36),
             font_size=10, color=GREEN if i == 4 else LIGHT_GRAY)

add_rect(s, Inches(9.35), Inches(6.8), Inches(3.5), Inches(0.35), INDIGO_SOFT)
add_text(s, "Schéma flexible — pas de migration",
         Inches(9.35), Inches(6.83), Inches(3.5), Inches(0.32),
         font_size=10, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Diagramme séquence Auth
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Diagramme de Séquence — Authentification",
            "Flux : Inscription sécurisée", badge="04 — SÉQUENCE")

actors = ["Utilisateur", "Frontend\n(Next.js)", "Backend\n(Express)", "PostgreSQL"]
ax_positions = [Inches(1.0), Inches(3.8), Inches(7.2), Inches(11.0)]
actor_colors = [INDIGO, PURPLE, GREEN, YELLOW]

for ax, label, col in zip(ax_positions, actors, actor_colors):
    add_rect(s, ax - Inches(0.6), Inches(2.05), Inches(1.2), Inches(0.55), col)
    add_text(s, label, ax - Inches(0.6), Inches(2.1),
             Inches(1.2), Inches(0.6),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, ax - Pt(1.5), Inches(2.6), Pt(3), Inches(4.5), col)

# Messages
msgs = [
    (Inches(1.0), Inches(3.8), Inches(2.85), "POST /register  (email, password, name)", WHITE, True),
    (Inches(3.8), Inches(7.2), Inches(4.0),  "SELECT * WHERE email=$1",                 LIGHT_GRAY, False),
    (Inches(7.2), Inches(11.0), Inches(4.6), "email non trouvé → null",                 GREEN, False),
    (Inches(7.2), Inches(3.8), Inches(4.0),  "bcrypt.hash(password, 10)",               PURPLE, True),
    (Inches(7.2), Inches(11.0), Inches(4.6), "INSERT INTO users → {id, email, name}",   LIGHT_GRAY, False),
    (Inches(11.0), Inches(7.2), Inches(4.6), "user.id",                                 GREEN, False),
    (Inches(7.2), Inches(3.8), Inches(4.0),  "jwt.sign({userId, email}, SECRET, 7d)",   INDIGO, True),
    (Inches(3.8), Inches(1.0), Inches(2.85), "201 · { token, user }",                   GREEN, False),
]
y_positions = [Inches(2.95), Inches(3.35), Inches(3.6), Inches(3.85),
               Inches(4.15), Inches(4.45), Inches(4.75), Inches(5.15)]

for (x1, x2, _, label, col, is_req), yy in zip(msgs, y_positions):
    # Ligne horizontale
    lx = min(x1, x2)
    lw = abs(x2 - x1)
    add_rect(s, lx, yy + Pt(8), lw, Pt(1.5), col)
    # Flèche
    arr_x = x2 - Pt(8) if x2 > x1 else x2 + Pt(3)
    add_text(s, "▶" if x2 > x1 else "◀", arr_x, yy, Pt(14), Pt(18),
             font_size=9, color=col, align=PP_ALIGN.LEFT)
    # Label
    mid_x = (x1 + x2) / 2 - Inches(0.8)
    add_text(s, label, mid_x, yy - Pt(12), Inches(2.5), Pt(18),
             font_size=9, color=col, align=PP_ALIGN.CENTER)

# Note finale
add_rect(s, Inches(0.55), Inches(5.8), Inches(12.3), Inches(0.55), INDIGO_SOFT)
add_text(s, "  Token JWT stocké en localStorage · Envoyé en Authorization: Bearer <token> à chaque requête protégée",
         Inches(0.65), Inches(5.85), Inches(12.0), Inches(0.45),
         font_size=11, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Démonstration
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, H, RGBColor(0x08, 0x10, 0x20))
add_rect(s, 0, 0, W, Inches(0.15), INDIGO)
add_rect(s, 0, Inches(7.35), W, Inches(0.15), INDIGO)

add_text(s, "🎬", Inches(5.9), Inches(0.8), Inches(1.5), Inches(1.5),
         font_size=54, align=PP_ALIGN.CENTER)
add_text(s, "DÉMONSTRATION LIVE",
         Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.2),
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Application My Smart Budget — 8 minutes",
         Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.7),
         font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

demo_steps = [
    ("01", "Login / Inscription",       "Interface sécurisée · validation JWT"),
    ("02", "Questionnaire 7 étapes",    "Animations · profil · conseils personnalisés"),
    ("03", "Dashboard & statistiques",  "Recharts · API /stats · données temps réel"),
    ("04", "Transactions — CRUD",       "Ajout · modification · suppression · IDOR"),
    ("05", "Profil & objectifs",        "GoalsSection · persistance PostgreSQL"),
]
for i, (num, title, sub) in enumerate(demo_steps):
    x = Inches(0.55) + i * Inches(2.55)
    add_rect(s, x, Inches(4.2), Inches(2.4), Inches(2.8), DARK_CARD)
    add_rect(s, x, Inches(4.2), Inches(2.4), Inches(0.5), INDIGO)
    add_text(s, num, x, Inches(4.22), Inches(2.4), Inches(0.46),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + Pt(8), Inches(4.82), Inches(2.2), Inches(0.5),
             font_size=11, bold=True, color=WHITE)
    add_text(s, sub,   x + Pt(8), Inches(5.3),  Inches(2.2), Inches(1.5),
             font_size=10, color=LIGHT_GRAY, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Code Frontend
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Code Frontend — Next.js 15",
            "Points techniques clés", badge="06 — FRONTEND")

# Code block 1
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(2.5), RGBColor(0x0A, 0x0F, 0x1E))
add_rect(s, Inches(0.55), Inches(2.05), Inches(6.0), Inches(0.4), RGBColor(0x16, 0x1C, 0x32))
add_text(s, "questionnaire.tsx — Validation par étape",
         Inches(0.7), Inches(2.08), Inches(5.8), Inches(0.36),
         font_size=10, bold=True, color=INDIGO)
code1 = [
    ("const canNext = () => {",                    WHITE),
    ("  if (step === 1) return !!answers.profileType;",    LIGHT_GRAY),
    ("  if (step === 4) return answers.dettes.length > 0;", LIGHT_GRAY),
    ("  if (step === 5) return answers.objectifsEpargne.length > 0;", LIGHT_GRAY),
    ("  return true;",                             GREEN),
    ("};",                                         WHITE),
]
for i, (line, col) in enumerate(code1):
    add_text(s, line, Inches(0.7), Inches(2.5) + i * Pt(23),
             Inches(5.7), Pt(22), font_size=10.5, color=col)

# Code block 2
add_rect(s, Inches(0.55), Inches(4.65), Inches(6.0), Inches(2.5), RGBColor(0x0A, 0x0F, 0x1E))
add_rect(s, Inches(0.55), Inches(4.65), Inches(6.0), Inches(0.4), RGBColor(0x16, 0x1C, 0x32))
add_text(s, "dashboard/page.tsx — Protection de route (client)",
         Inches(0.7), Inches(4.68), Inches(5.8), Inches(0.36),
         font_size=10, bold=True, color=INDIGO)
code2 = [
    ("useEffect(() => {",             WHITE),
    ("  const token = localStorage",  LIGHT_GRAY),
    ("        .getItem('token');",    LIGHT_GRAY),
    ("  if (!token) {",               WHITE),
    ("    router.push('/login');",    RED_SOFT),
    ("  }",                           WHITE),
    ("}, []);",                        WHITE),
]
for i, (line, col) in enumerate(code2):
    add_text(s, line, Inches(0.7), Inches(5.1) + i * Pt(23),
             Inches(5.7), Pt(22), font_size=10.5, color=col)

# Points techniques droite
points = [
    ("App Router Next.js",      "Routing basé fichiers · page.tsx · layout.tsx"),
    ("7 étapes animées",        "État local React · transition CSS conditionnelle"),
    ("Conseils dynamiques",     "generateAdvice() selon les réponses du profil"),
    ("Recharts",                "Graphiques interactifs sur les données de l'API"),
    ("TypeScript strict",       "Typage des props, états et réponses API"),
    ("Responsive Tailwind",     "Mobile first · classes conditionnelles · grid"),
    ("Profils différenciés",    "6 pages profil : étudiant, salarié, famille…"),
]
add_rect(s, Inches(6.75), Inches(2.05), Inches(6.1), Inches(5.1), DARK_CARD)
add_text(s, "Points clés", Inches(6.95), Inches(2.12), Inches(5.9), Inches(0.45),
         font_size=14, bold=True, color=INDIGO)
for i, (t, d) in enumerate(points):
    y = Inches(2.65) + i * Inches(0.63)
    add_rect(s, Inches(6.95), y, Pt(6), Inches(0.4), INDIGO)
    add_text(s, t,  Inches(7.2), y,          Inches(2.0), Inches(0.35),
             font_size=12, bold=True, color=WHITE)
    add_text(s, d,  Inches(9.25), y,         Inches(3.4), Inches(0.35),
             font_size=11, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Code Backend
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Code Backend — API REST Express",
            "Node.js · pg-promise · Mongoose", badge="07 — BACKEND")

# Table routes
add_rect(s, Inches(0.55), Inches(2.05), Inches(5.5), Inches(5.1), DARK_CARD)
add_rect(s, Inches(0.55), Inches(2.05), Inches(5.5), Inches(0.42), INDIGO_SOFT)
add_text(s, "  Méthode   Route                Description",
         Inches(0.6), Inches(2.08), Inches(5.4), Inches(0.36),
         font_size=10, bold=True, color=INDIGO)
routes = [
    ("POST",   "/api/auth/register",      "Inscription"),
    ("POST",   "/api/auth/login",         "Connexion JWT"),
    ("GET",    "/api/auth/verify",        "Vérif. token"),
    ("GET",    "/api/transactions",       "Liste transactions"),
    ("POST",   "/api/transactions/add",   "Créer transaction"),
    ("GET",    "/api/transactions/stats", "Stats mensuelles"),
    ("PUT",    "/api/transactions/:id",   "Modifier (ownerGuard)"),
    ("DELETE", "/api/transactions/:id",   "Supprimer (ownerGuard)"),
    ("POST",   "/api/profile/save",       "Upsert profil"),
    ("GET",    "/api/profile",            "Récupérer profil"),
    ("GET",    "/api/reports",            "Rapports MongoDB"),
]
method_colors = {
    "GET": GREEN, "POST": INDIGO, "PUT": YELLOW, "DELETE": RED_SOFT
}
for i, (meth, route, desc) in enumerate(routes):
    y = Inches(2.55) + i * Inches(0.43)
    bg_c = DARK_CARD if i % 2 == 0 else RGBColor(0x14, 0x1E, 0x35)
    add_rect(s, Inches(0.55), y, Inches(5.5), Inches(0.43), bg_c)
    col = method_colors.get(meth, WHITE)
    add_text(s, meth,  Inches(0.7), y + Pt(4), Inches(0.85), Pt(22),
             font_size=10, bold=True, color=col)
    add_text(s, route, Inches(1.6), y + Pt(4), Inches(2.5), Pt(22),
             font_size=10, color=LIGHT_GRAY)
    add_text(s, desc,  Inches(4.2), y + Pt(4), Inches(1.7), Pt(22),
             font_size=10, color=WHITE)

# Code SQL stats
add_rect(s, Inches(6.3), Inches(2.05), Inches(6.55), Inches(2.35), RGBColor(0x0A, 0x0F, 0x1E))
add_rect(s, Inches(6.3), Inches(2.05), Inches(6.55), Inches(0.4), RGBColor(0x16, 0x1C, 0x32))
add_text(s, "Agrégation SQL — Route /stats",
         Inches(6.45), Inches(2.08), Inches(6.3), Inches(0.36),
         font_size=10, bold=True, color=GREEN)
sql_lines = [
    ("SELECT",                                     INDIGO),
    ("  SUM(CASE WHEN amount > 0",                 LIGHT_GRAY),
    ("    THEN amount ELSE 0 END) AS total_revenus,", LIGHT_GRAY),
    ("  SUM(CASE WHEN amount < 0",                 LIGHT_GRAY),
    ("    THEN ABS(amount) ELSE 0 END) AS total_depenses,", LIGHT_GRAY),
    ("  COUNT(*) AS total_transactions",           GREEN),
    ("FROM transactions",                          WHITE),
    ("WHERE user_id=$1 AND date >= NOW()-'1 month'", YELLOW),
]
for i, (line, col) in enumerate(sql_lines):
    add_text(s, line, Inches(6.45), Inches(2.5) + i * Pt(20),
             Inches(6.2), Pt(20), font_size=10, color=col)

# Code Upsert
add_rect(s, Inches(6.3), Inches(4.5), Inches(6.55), Inches(2.4), RGBColor(0x0A, 0x0F, 0x1E))
add_rect(s, Inches(6.3), Inches(4.5), Inches(6.55), Inches(0.4), RGBColor(0x16, 0x1C, 0x32))
add_text(s, "Upsert PostgreSQL — Route /profile/save",
         Inches(6.45), Inches(4.53), Inches(6.3), Inches(0.36),
         font_size=10, bold=True, color=PURPLE)
upsert = [
    ("INSERT INTO user_profiles",               INDIGO),
    ("  (user_id, profile_type, ...)",          LIGHT_GRAY),
    ("VALUES ($1, $2, ...)",                    LIGHT_GRAY),
    ("ON CONFLICT (user_id) DO UPDATE SET",     YELLOW),
    ("  profile_type = $2,",                    GREEN),
    ("  revenus_mensuels = $3 ...",             GREEN),
    ("RETURNING *",                              WHITE),
]
for i, (line, col) in enumerate(upsert):
    add_text(s, line, Inches(6.45), Inches(4.95) + i * Pt(21),
             Inches(6.2), Pt(21), font_size=10, color=col)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Sécurité
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Sécurité Applicative",
            "OWASP — Bonnes pratiques implementées", badge="08 — SÉCURITÉ")

sec_items = [
    (INDIGO,  "JWT + bcrypt",
     "Authentification stateless · Tokens 7j\nbcrypt salt=10 · Mot de passe validé (regex 8+ car, maj, min, chiffre)"),
    (PURPLE,  "Helmet",
     "En-têtes HTTP sécurisés automatiques\nContent-Security-Policy · X-Frame-Options · X-Content-Type-Options\nX-Powered-By désactivé"),
    (RED_SOFT,"Rate Limiting",
     "20 requêtes max / 15 min sur /api/auth\nProtection brute force sur login et register\nRejet avec message explicite"),
    (GREEN,   "Protection IDOR (ownerGuard)",
     "Vérifie que la transaction appartient\nà l'utilisateur extrait du JWT\nAucune attaque par manipulation d'id possible"),
    (YELLOW,  "userId issu du JWT",
     "Jamais accepté depuis le body de la requête\nExtrait exclusivement du token décodé\nUn user ne peut accéder qu'à ses propres données"),
    (INDIGO,  "CORS + Payload limité",
     "Origines autorisées explicitement configurées\nBody JSON limité à 10kb\nMéthodes HTTP restreintes : GET POST PUT DELETE"),
]

for i, (col, title, desc) in enumerate(sec_items):
    x = Inches(0.55) + (i % 3) * Inches(4.25)
    y = Inches(2.05) + (i // 3) * Inches(2.55)
    add_rect(s, x, y, Inches(4.1), Inches(2.38), DARK_CARD)
    add_rect(s, x, y, Inches(4.1), Inches(0.45), col)
    add_text(s, title, x + Pt(8), y + Pt(4), Inches(3.9), Inches(0.38),
             font_size=13, bold=True, color=BG_DARK)
    add_text(s, desc, x + Pt(8), y + Inches(0.55), Inches(3.9), Inches(1.75),
             font_size=11, color=LIGHT_GRAY, wrap=True)

# Code ownerGuard highlight
add_rect(s, Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.38), INDIGO_SOFT)
add_text(s, "  ownerGuard.js : SELECT id FROM transactions WHERE id=$1 AND user_id=$2  →  protection IDOR",
         Inches(0.65), Inches(7.08), Inches(12.0), Inches(0.32),
         font_size=11, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Tests
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Tests & Qualité Logicielle",
            "Jest · Supertest · SonarCloud", badge="09 — TESTS")

# Pyramide
add_rect(s, Inches(0.55), Inches(2.05), Inches(5.7), Inches(5.15), DARK_CARD)
add_text(s, "Pyramide des tests", Inches(0.75), Inches(2.12), Inches(5.4), Inches(0.45),
         font_size=14, bold=True, color=INDIGO)

pyramid = [
    (GREEN,  "Tests E2E",           "Parcours utilisateur complet",        Inches(1.8),  Inches(2.8)),
    (YELLOW, "Tests Intégration",   "Routes HTTP · Supertest · DB mockée", Inches(1.3),  Inches(2.95)),
    (INDIGO, "Tests Unitaires",     "Middleware auth · Fonctions pures",   Inches(0.9),  Inches(2.95)),
]
tri_data = [
    (Inches(3.1), Inches(2.65), Inches(1.6), GREEN),
    (Inches(2.6), Inches(3.55), Inches(2.6), YELLOW),
    (Inches(2.15), Inches(4.3), Inches(3.45), INDIGO),
]
for (bx, by, bw, col), (_, label, _, _, _) in zip(tri_data, pyramid):
    add_rect(s, bx, by, bw, Inches(0.85), col)

add_text(s, "E2E\n(prévu)", Inches(3.1), Inches(2.7), Inches(1.6), Inches(0.75),
         font_size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
add_text(s, "INTÉGRATION\nJest + Supertest", Inches(2.6), Inches(3.6), Inches(2.6), Inches(0.75),
         font_size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
add_text(s, "UNITAIRES — auth.middleware.test.js", Inches(2.15), Inches(4.35), Inches(3.45), Inches(0.75),
         font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Liste des tests
tests = [
    (GREEN, "PASS", "Register — champs manquants → 400"),
    (GREEN, "PASS", "Register — mot de passe faible → 400"),
    (GREEN, "PASS", "Register — email déjà utilisé → 400"),
    (GREEN, "PASS", "Register — succès → 201 + token"),
    (GREEN, "PASS", "Login — champs manquants → 400"),
    (GREEN, "PASS", "Login — email inexistant → 401"),
    (GREEN, "PASS", "Login — mauvais mot de passe → 401"),
    (GREEN, "PASS", "Login — succès → 200 + token"),
    (GREEN, "PASS", "Verify — token manquant → 401"),
    (GREEN, "PASS", "Verify — token valide → valid:true"),
]
add_rect(s, Inches(6.5), Inches(2.05), Inches(6.35), Inches(5.15), DARK_CARD)
add_text(s, "Tests auth.routes.test.js — 10 / 10", Inches(6.7), Inches(2.12),
         Inches(6.1), Inches(0.45), font_size=14, bold=True, color=GREEN)
for i, (col, status, label) in enumerate(tests):
    y = Inches(2.65) + i * Inches(0.42)
    add_rect(s, Inches(6.7), y, Inches(0.55), Inches(0.32), col)
    add_text(s, status, Inches(6.7), y + Pt(2), Inches(0.55), Pt(24),
             font_size=9, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s, label, Inches(7.35), y + Pt(2), Inches(5.3), Pt(24),
             font_size=11, color=WHITE)

# SonarCloud note
add_rect(s, Inches(6.5), Inches(6.85), Inches(6.35), Inches(0.35), INDIGO_SOFT)
add_text(s, "  SonarCloud · coverage lcov · sonar-project.properties · exclusion node_modules & .next",
         Inches(6.6), Inches(6.88), Inches(6.15), Inches(0.3),
         font_size=10, color=INDIGO)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Déploiement
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Déploiement & DevOps",
            "Docker Compose · 4 services · Docker Hub", badge="10 — DEVOPS")

services = [
    (INDIGO,  "frontend",  "Next.js 15",      ":3000", "abedel/my-smart-budget-frontend"),
    (PURPLE,  "backend",   "Node.js Express", ":5001", "abedel/my-smart-budget-backend"),
    (GREEN,   "postgres",  "PostgreSQL 16",   ":5432", "postgres:16-alpine"),
    (YELLOW,  "mongodb",   "MongoDB 7",       ":27017","mongo:7"),
]

for i, (col, name, tech, port, image) in enumerate(services):
    x = Inches(0.55) + i * Inches(3.2)
    y = Inches(2.05)
    w = Inches(3.05)
    add_rect(s, x, y, w, Inches(3.5), DARK_CARD)
    add_rect(s, x, y, w, Inches(0.5), col)
    add_text(s, name, x, y + Pt(4), w, Inches(0.42),
             font_size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s, tech,  x + Pt(8), y + Inches(0.62), w - Pt(12), Inches(0.42),
             font_size=12, color=WHITE, bold=True)
    add_text(s, f"Port : {port}", x + Pt(8), y + Inches(1.05), w - Pt(12), Inches(0.35),
             font_size=11, color=LIGHT_GRAY)
    add_text(s, f"Image :\n{image}", x + Pt(8), y + Inches(1.45), w - Pt(12), Inches(0.8),
             font_size=10, color=LIGHT_GRAY)
    if name in ("postgres", "mongodb"):
        add_rect(s, x + Pt(8), y + Inches(2.6), w - Pt(16), Inches(0.4), INDIGO_SOFT)
        add_text(s, "Healthcheck actif", x + Pt(8), y + Inches(2.62), w - Pt(16), Inches(0.36),
                 font_size=10, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

# Réseau
add_rect(s, Inches(0.55), Inches(5.75), Inches(12.3), Inches(0.6), DARK_CARD)
add_text(s, "  🌐  Réseau interne Docker :  smart-budget-network (bridge) · restart: always · depends_on: service_healthy",
         Inches(0.7), Inches(5.82), Inches(12.0), Inches(0.46),
         font_size=12, color=WHITE)

# prod info
add_rect(s, Inches(0.55), Inches(6.5), Inches(12.3), Inches(0.55), INDIGO_SOFT)
add_text(s, "  Production : https://smarterbudget.net  —  NEXT_PUBLIC_API_URL injectée via variable d'environnement",
         Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.45),
         font_size=12, bold=True, color=INDIGO)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Bilan & Conclusion",
            "Compétences CDA couvertes", badge="CONCLUSION")

ccps = [
    (INDIGO, "CCP1 — Développer une application sécurisée",
     ["Interface utilisateur Next.js complète et responsive",
      "Composants backend Express avec API REST",
      "Gestion de projet Agile avec Git et SonarCloud"]),
    (PURPLE, "CCP2 — Concevoir et développer en couches",
     ["Architecture 3-tiers strictement respectée",
      "BDD relationnelle PostgreSQL + NoSQL MongoDB",
      "Modélisation MCD/MLD · Use Cases · Séquences"]),
    (GREEN, "CCP3 — Préparer le déploiement",
     ["Tests automatisés Jest + Supertest (10 tests, 100%)",
      "Déploiement Docker Compose · Docker Hub",
      "Qualité continue SonarCloud · couverture lcov"]),
]

for i, (col, title, items) in enumerate(ccps):
    x = Inches(0.55) + i * Inches(4.25)
    add_rect(s, x, Inches(2.05), Inches(4.1), Inches(3.2), DARK_CARD)
    add_rect(s, x, Inches(2.05), Inches(4.1), Inches(0.52), col)
    add_text(s, title, x + Pt(8), Inches(2.08), Inches(3.9), Inches(0.48),
             font_size=11, bold=True, color=BG_DARK, wrap=True)
    for j, item in enumerate(items):
        add_rect(s, x + Pt(8), Inches(2.67) + j * Inches(0.8), Pt(7), Inches(0.52), col)
        add_text(s, item, x + Pt(22), Inches(2.65) + j * Inches(0.8),
                 Inches(3.7), Inches(0.7), font_size=12, color=LIGHT_GRAY, wrap=True)

# Apprentissages
add_rect(s, Inches(0.55), Inches(5.45), Inches(12.3), Inches(1.8), DARK_CARD)
add_text(s, "Points forts & apprentissages clés",
         Inches(0.75), Inches(5.5), Inches(12.0), Inches(0.42),
         font_size=14, bold=True, color=YELLOW)
learnings = [
    "Polyglot persistence (PostgreSQL + MongoDB) selon les besoins métier",
    "Protection IDOR avec ownerGuard — sécurité souvent négligée",
    "userId extrait du JWT, jamais du body — principe de non-répudiation",
    "Fail-fast au démarrage si variables d'environnement manquantes",
]
for i, l in enumerate(learnings):
    xi = Inches(0.75) if i < 2 else Inches(6.75)
    yi = Inches(6.0) + (i % 2) * Inches(0.55)
    add_text(s, f"▸  {l}", xi, yi, Inches(5.8), Inches(0.45),
             font_size=12, color=LIGHT_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Questions
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, H, RGBColor(0x08, 0x10, 0x20))
add_rect(s, 0, 0, W, Inches(0.15), INDIGO)
add_rect(s, 0, Inches(7.35), W, Inches(0.15), INDIGO)

add_text(s, "?", Inches(5.9), Inches(0.6), Inches(1.5), Inches(1.5),
         font_size=72, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

add_text(s, "Questions",
         Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.2),
         font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Je suis disponible pour répondre à toutes vos questions",
         Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.6),
         font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

qa = [
    "Architecture 3-tiers",
    "Protection IDOR / Sécurité",
    "Choix PostgreSQL + MongoDB",
    "Déploiement Docker",
    "Tests Jest / Supertest",
    "Questionnaire & profils",
]
for i, q in enumerate(qa):
    x = Inches(1.5) + (i % 3) * Inches(3.5)
    y = Inches(4.0) + (i // 3) * Inches(0.85)
    add_rect(s, x, y, Inches(3.3), Inches(0.65), INDIGO_SOFT)
    add_text(s, q, x, y + Pt(7), Inches(3.3), Inches(0.5),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "My Smart Budget — Concepteur Développeur d'Applications | Colint School 2025",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
         font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ─── Sauvegarde ──────────────────────────────────────────────────────────────
output = "/Users/ahbeuhdellsaidi/Downloads/cda-folder-main/dssier_pfe/Soutenance_CDA_MySmartBudget.pptx"
prs.save(output)
print(f"✅ PowerPoint généré : {output}")
print(f"   {len(prs.slides)} slides")
