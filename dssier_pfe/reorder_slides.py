"""
Restructure the 50-slide PPTX:
- Slides 1-14  (indices 0-13)  : originals Title → Déploiement  → KEEP in place
- Slides 15-16 (indices 14-15) : Conclusion + Questions           → MOVE to end
- Slides 17-50 (indices 16-49) : New detail slides                → INSERT after slide 14
- Slide 49 of new set (s50 "Bilan & Perspectives")                → rename to "Roadmap"

Final order: [0-13, 16-49, 14, 15]  →  50 slides
  Slides 1-14  : original compact slides
  Slides 15-48 : new detailed slides (Conception → DevOps)
  Slide 49     : Roadmap & Perspectives  (was my s50)
  Slide 50     : original Conclusion CCP
  → then Questions removed (stays as slide 50 if we keep it)

Wait — we want: 1-14 original, 15-48 new detail, 49 = original Conclusion, 50 = Questions
So we drop my s50 "Bilan & Perspectives" (index 49) and keep 49 slides...
No — we keep all 50 but rename my slide to "Perspectives" (not Conclusion).

Order: [0,1,2,3,4,5,6,7,8,9,10,11,12,13,  16,17,18,...,49,  14,15]
= 14 original + 34 new + Conclusion + Questions = 50
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree

PPTX = "/Users/ahbeuhdellsaidi/Downloads/cda-folder-main/dssier_pfe/soutenance_CDA_MySmartBudget.pptx"
prs = Presentation(PPTX)

# ── Step 1: rename my "CONCLUSION / Bilan & Perspectives" slide (index 49)
#            to "PERSPECTIVES / Roadmap & Perspectives" so it doesn't clash
#            with the original Conclusion that will follow it.
slide_49 = prs.slides[49]
for shape in slide_49.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text == "CONCLUSION":
                run.text = "PERSPECTIVES"
            elif run.text == "Bilan & Perspectives":
                run.text = "Roadmap & Perspectives"
            elif run.text == "My Smart Budget — du concept à l'application déployable":
                run.text = "Les prochaines étapes naturelles du projet"

# ── Step 2: reorder slides via XML manipulation
# New order (0-based indices): 0-13, 16-49, 14, 15
new_order = list(range(14)) + list(range(16, 50)) + [14, 15]

# python-pptx stores slide references in prs.slides._sldIdLst
xml_slides = prs.slides._sldIdLst
# collect current <p:sldId> elements in order
sld_id_elems = list(xml_slides)

assert len(sld_id_elems) == 50, f"Expected 50 slides, got {len(sld_id_elems)}"

# Build reordered list
reordered = [sld_id_elems[i] for i in new_order]

# Clear sldIdLst and re-append in new order
for elem in sld_id_elems:
    xml_slides.remove(elem)
for elem in reordered:
    xml_slides.append(elem)

# ── Step 3: verify
assert len(prs.slides) == 50
print("Slide order after restructure:")
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.paragraphs[0].text.strip()
            if t:
                texts.append(t)
    print(f"  {i+1:2d}. {texts[1] if len(texts) > 1 else texts[0] if texts else '(empty)'}")

prs.save(PPTX)
print(f"\nSaved — {len(prs.slides)} slides total.")
